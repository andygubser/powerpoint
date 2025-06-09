import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import random

# ================================
# CONFIGURATION PARAMETERS
# ================================

# File settings
INPUT_FILE = "../data/words.xlsx"  # Excel file with words (relative to src folder)
SHEET_NAME = "Sheet1"
OUTPUT_FILE = "../output/words_presentation_3sec.pptx"  # Output location

# Font settings
FONT_NAME = "DCH-Basisschrift"  # Font to use (install from ../fonts/ folder)
MAX_FONT_SIZE = 320  # Maximum font size in points
MIN_FONT_SIZE = 20  # Minimum font size in points
FONT_STEP_SIZE = 5  # How much to decrease font size per step
FONT_BOLD = True  # Whether text should be bold
FONT_COLOR_RGB = (0, 0, 0)  # Font color as RGB tuple (Black = 0,0,0)

# Slide layout settings
SLIDE_MARGIN_LEFT = 0.75  # Left margin in inches
SLIDE_MARGIN_TOP = 1.5  # Top margin in inches
SLIDE_WIDTH = 8.5  # Text area width in inches
SLIDE_HEIGHT = 5  # Text area height in inches
MAX_TEXT_WIDTH = 8.5  # Maximum text width for fitting calculation

# ⏱️ TIMING SETTINGS - EASILY ADJUSTABLE
AUTO_ADVANCE_SECONDS = (
    3  # 🔧 CHANGE THIS: How many seconds per slide (e.g., 1, 2, 3, 5, 10)
)
DISABLE_MOUSE_CLICK = True  # Whether to disable mouse click advancement

# Presentation settings
RANDOMIZE_ORDER = True  # Whether to shuffle words randomly

# Character width estimation factor (adjust if needed)
CHAR_WIDTH_FACTOR = 0.6  # Approximation: each char takes this * font_size in width

# Progress reporting
SHOW_PROGRESS_EVERY = 20  # Show progress every N slides
SHOW_FIRST_N_SLIDES = 10  # Always show progress for first N slides

# ================================
# END CONFIGURATION
# ================================


def find_optimal_font_size(
    word,
    max_width_inches=MAX_TEXT_WIDTH,
    max_font_size=MAX_FONT_SIZE,
    min_font_size=MIN_FONT_SIZE,
):
    """
    Calculate optimal font size based on word length and available space
    Approximates text width using character count and font size
    """
    # Convert inches to points (1 inch = 72 points)
    max_width_points = max_width_inches * 72

    for font_size in range(max_font_size, min_font_size - 1, -FONT_STEP_SIZE):
        # Estimate text width: char_count * font_size * factor
        estimated_width = len(word) * font_size * CHAR_WIDTH_FACTOR

        if estimated_width <= max_width_points:
            return font_size

    return min_font_size


def set_slide_transitions(presentation, advance_time_seconds=3, click_advance=False):
    """
    Set slide transition properties for all slides using XML manipulation
    Clears all existing timings first, then sets clean auto-advance timing
    """
    try:
        from lxml import etree

        # Convert seconds to milliseconds for XML
        advance_time_ms = advance_time_seconds * 1000

        transitions_set = 0

        for slide in presentation.slides:
            slide_element = slide.element

            # Remove any existing transition elements first (reset to null)
            existing_transitions = slide_element.findall(
                ".//{http://schemas.openxmlformats.org/presentationml/2006/main}transition"
            )
            for transition in existing_transitions:
                slide_element.remove(transition)

            # Create fresh transition element
            transition_elem = etree.SubElement(
                slide_element,
                "{http://schemas.openxmlformats.org/presentationml/2006/main}transition",
            )

            # Set transition properties (clean slate)
            transition_elem.set("spd", "fast")  # Transition speed

            # Disable mouse click advancement
            if click_advance:
                transition_elem.set("advanceOnClick", "1")
            else:
                transition_elem.set("advanceOnClick", "0")  # Explicitly disable clicks

            # Set automatic advance timing (3 seconds)
            transition_elem.set("advanceAfterTime", str(advance_time_ms))

            # Add additional attributes for better compatibility
            transition_elem.set("dur", str(advance_time_ms))  # Duration

            transitions_set += 1

        print(f"✅ Reset and configured transitions on {transitions_set} slides")
        print(
            f"   • Duration reset to null, then set to {advance_time_seconds}s auto-advance"
        )
        print(f"   • Mouse click: {'Enabled' if click_advance else 'Disabled'}")
        return True

    except ImportError:
        print("⚠️  lxml not available - cannot set transitions programmatically")
        return False
    except Exception as e:
        print(f"⚠️  Error setting transitions: {str(e)[:50]}...")
        return False


def create_word_presentation():
    print("=" * 60)
    print("POWERPOINT WORD PRESENTATION GENERATOR")
    print("=" * 60)
    print(f"📁 Input file: {INPUT_FILE}")
    print(f"📄 Output file: {OUTPUT_FILE}")
    print(f"🎨 Font: {FONT_NAME} ({MIN_FONT_SIZE}pt - {MAX_FONT_SIZE}pt)")
    print(f"🎲 Randomize order: {RANDOMIZE_ORDER}")
    print(f"⏱️  Auto-advance: {AUTO_ADVANCE_SECONDS} seconds (automatic)")
    print(f"🖱️  Mouse click: {'Disabled' if DISABLE_MOUSE_CLICK else 'Enabled'}")
    print("-" * 60)

    # Read the Excel file without treating first row as header
    df = pd.read_excel(INPUT_FILE, sheet_name=SHEET_NAME, header=None)

    # Get all words from the first column (index 0)
    words = df[0].dropna().tolist()

    # Randomize the order of words if enabled
    if RANDOMIZE_ORDER:
        random.shuffle(words)
        print(f"🔀 Words randomized. First few: {words[:5]}")
    else:
        print(f"📝 Words in original order. First few: {words[:5]}")

    print(f"📊 Creating presentation with {len(words)} words...")

    # Create a new presentation
    prs = Presentation()

    # Define slide layout (using blank layout)
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout

    font_adjustments = 0
    size_stats = {"max_size": 0, "large": 0, "medium": 0, "small": 0}

    for i, word in enumerate(words):
        # Add a slide
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add a text box with configured dimensions
        left = Inches(SLIDE_MARGIN_LEFT)
        top = Inches(SLIDE_MARGIN_TOP)
        width = Inches(SLIDE_WIDTH)
        height = Inches(SLIDE_HEIGHT)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        # Configure text frame for perfect centering
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center vertically
        text_frame.word_wrap = False  # Force single line
        text_frame.auto_size = None  # Manual control
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0

        # Configure the text
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = str(word)
        paragraph.alignment = PP_ALIGN.CENTER  # Center horizontally

        # Calculate optimal font size for this word
        optimal_size = find_optimal_font_size(word)

        # Format the text with configured font settings
        font = paragraph.font
        font.name = FONT_NAME
        font.size = Pt(optimal_size)
        font.color.rgb = RGBColor(*FONT_COLOR_RGB)
        font.bold = FONT_BOLD

        # Track statistics
        if optimal_size == MAX_FONT_SIZE:
            size_stats["max_size"] += 1
        elif optimal_size >= MAX_FONT_SIZE * 0.75:  # 75% of max
            size_stats["large"] += 1
        elif optimal_size >= MAX_FONT_SIZE * 0.5:  # 50% of max
            size_stats["medium"] += 1
        else:
            size_stats["small"] += 1

        if optimal_size < MAX_FONT_SIZE:
            font_adjustments += 1

        # Progress reporting
        if i < SHOW_FIRST_N_SLIDES or i % SHOW_PROGRESS_EVERY == 0:
            if optimal_size < MAX_FONT_SIZE:
                print(f"  📄 Slide {i+1}: '{word}' (font: {optimal_size}pt)")
            else:
                print(f"  📄 Slide {i+1}: '{word}'")

    # Set slide transitions programmatically
    print(f"\n⚙️  Configuring slide transitions...")
    print(f"   • Resetting all existing slide durations to null")
    print(f"   • Setting auto-advance to {AUTO_ADVANCE_SECONDS} seconds")
    print(
        f"   • Mouse click advancement: {'Disabled' if DISABLE_MOUSE_CLICK else 'Enabled'}"
    )

    transition_success = set_slide_transitions(
        prs,
        advance_time_seconds=AUTO_ADVANCE_SECONDS,
        click_advance=not DISABLE_MOUSE_CLICK,
    )

    # Save the presentation
    prs.save(OUTPUT_FILE)

    print("\n" + "=" * 60)
    print("✅ PRESENTATION CREATED SUCCESSFULLY!")
    print("=" * 60)
    print(f"📁 File saved: {OUTPUT_FILE}")
    print(f"📊 Total slides: {len(words)}")
    print(f"🔧 Font adjustments: {font_adjustments} slides")
    print(f"\n📈 Size distribution:")
    print(f"  • Maximum ({MAX_FONT_SIZE}pt): {size_stats['max_size']} words")
    print(f"  • Large (≥{int(MAX_FONT_SIZE*0.75)}pt): {size_stats['large']} words")
    print(f"  • Medium (≥{int(MAX_FONT_SIZE*0.5)}pt): {size_stats['medium']} words")
    print(f"  • Small (<{int(MAX_FONT_SIZE*0.5)}pt): {size_stats['small']} words")

    print(f"\n✨ Features enabled:")
    print(f"  ✓ Perfect centering (horizontal & vertical)")
    print(f"  ✓ Auto-fit sizing (no overflow or line breaks)")
    print(f"  ✓ Custom font: {FONT_NAME}")
    if RANDOMIZE_ORDER:
        print(f"  ✓ Random order (zufällig)")
    else:
        print(f"  ✓ Original order maintained")
    print(f"  ✓ {AUTO_ADVANCE_SECONDS}-second auto-advance")
    print(f"  ✓ Mouse click {'disabled' if DISABLE_MOUSE_CLICK else 'enabled'}")

    if transition_success:
        print(f"\n🎉 READY TO USE:")
        print("  • Open the presentation and start slideshow")
        print(f"  • Slides auto-advance every {AUTO_ADVANCE_SECONDS} seconds")
        if DISABLE_MOUSE_CLICK:
            print("  • Mouse clicks are disabled")
        print("  • No manual setup required!")
    else:
        print(f"\n⏱️  MANUAL SETUP REQUIRED:")
        print("  1. Open the presentation in PowerPoint")
        print("  2. Go to 'Transitions' tab")
        if DISABLE_MOUSE_CLICK:
            print("  3. Uncheck 'On Mouse Click'")
        print(f"  4. Check 'After' and set to 00:0{AUTO_ADVANCE_SECONDS}")
        print("  5. Click 'Apply To All'")

    print(f"\n💡 To customize: Edit parameters at top of {__file__}")
    print(f"🎨 Font note: {FONT_NAME} must be installed for proper display")
    print("=" * 60)


if __name__ == "__main__":
    create_word_presentation()
